"""
Document Extractor - Quick text extraction from documents

Extracts text content from PDF, DOCX, PPTX files for immediate RAG queries.
Does NOT store in database - only provides context for current query.

Usage:
    from src.document_extractor import DocumentExtractor
    extractor = DocumentExtractor()
    text = extractor.extract_from_file(file_path)
"""

import os
import zipfile
import io
from typing import Optional, Dict, List, Any
from pathlib import Path


class DocumentExtractor:
    """
    Quick text extraction from documents for chat queries.
    
    Supports:
    - PDF files (via PyPDF)
    - DOCX files (via python-docx)
    - PPTX files (via python-pptx)
    """
    
    def __init__(self, max_chars: int = 50000):
        """
        Initialize document extractor.
        
        Args:
            max_chars: Maximum characters to extract (prevent huge context)
        """
        self.max_chars = max_chars
        print(f"✓ Document Extractor initialized (max: {max_chars} chars)")
    
    def extract_from_file(self, file_path: str) -> Dict[str, Any]:
        """
        Extract text and images from any supported document type.
        
        Args:
            file_path: Path to document file
            
        Returns:
            Dict containing:
            - 'text': Extracted text content
            - 'images': List of dicts {'name': str, 'bytes': bytes, 'type': str}
        """
        ext = os.path.splitext(file_path)[1].lower()
        
        print(f"\n📄 Extracting content from: {Path(file_path).name}")
        
        result = {
            "text": "",
            "images": []
        }
        
        try:
            if ext == '.pdf':
                result["text"] = self._extract_pdf(file_path)
            elif ext == '.docx':
                result = self._extract_docx(file_path)
            elif ext == '.pptx':
                result = self._extract_pptx(file_path)
            else:
                result["text"] = f"[Unsupported file type: {ext}]"
            
            # Truncate text if too long
            text = result.get("text", "")
            if len(text) > self.max_chars:
                result["text"] = text[:self.max_chars] + f"\n\n[... truncated, original length: {len(text)} chars]"
                print(f"⚠️  Text truncated to {self.max_chars} chars")
            
            print(f"✓ Extracted {len(result['text'])} characters")
            if result['images']:
                print(f"✓ Extracted {len(result['images'])} images")
                
            return result
            
        except Exception as e:
            print(f"❌ Error extracting file: {e}")
            result["text"] = f"[Error: {str(e)}]"
            return result
    
    def _extract_pdf(self, file_path: str) -> str:
        """Extract text from PDF using PyPDF."""
        try:
            from pypdf import PdfReader
            
            reader = PdfReader(file_path)
            text_parts = []
            
            # Extract from all pages
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(f"--- Page {i+1} ---\n{page_text}")
                
                # Stop if exceeding max chars
                current_length = sum(len(p) for p in text_parts)
                if current_length > self.max_chars:
                    break
            
            return "\n\n".join(text_parts)
            
        except Exception as e:
            return f"[Error extracting PDF: {str(e)}]"
    
    def _extract_docx(self, file_path: str) -> Dict[str, Any]:
        """Extract text and images from DOCX."""
        result = {"text": "", "images": []}
        try:
            from docx import Document
            
            doc = Document(file_path)
            text_parts = []
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
                
                # Stop if exceeding max chars
                current_length = sum(len(p) for p in text_parts)
                if current_length > self.max_chars:
                    break
            
            # Extract tables
            if sum(len(p) for p in text_parts) < self.max_chars:
                for table in doc.tables:
                    table_text = []
                    for row in table.rows:
                        # Handle potential merged cells or empty cells
                        row_cells = [cell.text.strip() for cell in row.cells] if row.cells else []
                        row_text = " | ".join(filter(None, row_cells))
                        if row_text:
                            table_text.append(f"| {row_text} |")
                    
                    if table_text:
                        text_parts.append("\n" + "\n".join(table_text) + "\n")
                    
                    if sum(len(p) for p in text_parts) > self.max_chars:
                        break
            
            # --- Extract Images via ZipFile ---
            # DOCX is a zip file. Images can be in 'media/' or 'word/media/'
            try:
                with zipfile.ZipFile(file_path) as z:
                    for filename in z.namelist():
                        # Check both 'media/' and 'word/media/' paths
                        if (filename.startswith('word/media/') or filename.startswith('media/')) and not filename.endswith('/'):
                            # Only extract valid image types
                            ext = os.path.splitext(filename)[1].lower()
                            if ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff']:
                                idx = len(result["images"])
                                short_name = Path(filename).name # image1.png
                                
                                # Read image bytes
                                img_bytes = z.read(filename)
                                
                                # Add to results
                                result["images"].append({
                                    "index": idx,
                                    "name": short_name,
                                    "bytes": img_bytes,
                                    "ext": ext
                                })
                                
                                # Append placeholder to text (since we don't know exact position)
                                # Use format: ![name](name) to match batch_process_all.py injection pattern
                                text_parts.append(f"\n![{short_name}]({short_name})\n")
            except Exception as img_err:
                print(f"⚠️  Error extracting DOCX images: {img_err}")
                
            result["text"] = "\n\n".join(text_parts)
            return result
            
        except Exception as e:
            return {"text": f"[Error extracting DOCX: {str(e)}]", "images": []}
    
    def _extract_pptx(self, file_path: str) -> Dict[str, Any]:
        """Extract text and images from PPTX using python-pptx."""
        result = {"text": "", "images": []}
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            
            prs = Presentation(file_path)
            text_parts = []
            
            # Extract from all slides
            for i, slide in enumerate(prs.slides):
                slide_header = f"--- Slide {i+1} ---"
                slide_content = [slide_header]
                
                # Iterate shapes
                for shape in slide.shapes:
                    # Text extraction
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text)
                    
                    # Image extraction
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image = shape.image
                            img_bytes = image.blob
                            ext = "." + image.ext
                            
                            idx = len(result["images"])
                            img_name = f"slide{i+1}_img{idx}{ext}"
                            
                            # Add to results
                            result["images"].append({
                                "index": idx,
                                "name": img_name,
                                "bytes": img_bytes,
                                "ext": ext,
                                "page": i  # Track slide number
                            })
                            
                            # Insert placeholder in text flow
                            # Use format: ![name](name) to match batch_process_all.py injection pattern
                            slide_content.append(f"\n![{img_name}]({img_name})\n")
                            
                        except Exception as img_err:
                            print(f"⚠️  Error extracting PPTX image shape: {img_err}")
                
                text_parts.append("\n".join(slide_content))
                
                # Stop if exceeding max chars
                current_length = sum(len(p) for p in text_parts)
                if current_length > self.max_chars:
                    break
            
            result["text"] = "\n\n".join(text_parts)
            return result
            
        except Exception as e:
            return {"text": f"[Error extracting PPTX: {str(e)}]", "images": []}


def get_file_type(filename: str) -> str:
    """
    Determine file type from filename.
    
    Returns: 'image', 'pdf', 'docx', 'pptx', or 'unknown'
    """
    ext = os.path.splitext(filename)[1].lower()
    
    if ext in ['.jpg', '.jpeg', '.png', '.gif', '.webp', '.bmp']:
        return 'image'
    elif ext == '.pdf':
        return 'pdf'
    elif ext == '.docx':
        return 'docx'
    elif ext == '.pptx':
        return 'pptx'
    else:
        return 'unknown'


# For testing
if __name__ == "__main__":
    import sys
    
    if len(sys.argv) < 2:
        print("Usage: python document_extractor.py <file_path>")
        sys.exit(1)
    
    file_path = sys.argv[1]
    
    extractor = DocumentExtractor()
    text = extractor.extract_from_file(file_path)
    
    print("\n" + "="*60)
    print("EXTRACTED TEXT")
    print("="*60)
    print(text)
    print("="*60)
